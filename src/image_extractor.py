"""Extracción y persistencia de imágenes desde PDF y PPTX."""
import io
from pathlib import Path
from typing import Optional

from PIL import Image
from loguru import logger


class ImageExtractor:
    def __init__(self, output_dir: Path) -> None:
        self.images_dir = output_dir / "images"
        self.images_dir.mkdir(parents=True, exist_ok=True)
        self._counter = 0

    def _next_path(self) -> Path:
        self._counter += 1
        return self.images_dir / f"image_{self._counter:03d}.png"

    def save_image(self, image: Image.Image) -> Optional[Path]:
        """Guarda una imagen en el directorio de salida y retorna su ruta."""
        try:
            path = self._next_path()
            image.convert("RGB").save(str(path), "PNG")
            return path
        except Exception as exc:
            logger.warning(f"No se pudo guardar imagen: {exc}")
            return None

    def extract_from_pdf(self, pdf_path: Path) -> list[Path]:
        """Extrae todas las imágenes embebidas en un PDF."""
        extracted: list[Path] = []
        try:
            import fitz

            doc = fitz.open(str(pdf_path))
            for page_num in range(len(doc)):
                for img_ref in doc[page_num].get_images(full=True):
                    xref = img_ref[0]
                    try:
                        base = doc.extract_image(xref)
                        img = Image.open(io.BytesIO(base["image"]))
                        path = self.save_image(img)
                        if path:
                            extracted.append(path)
                    except Exception as exc:
                        logger.debug(f"Imagen xref={xref} en página {page_num}: {exc}")
            doc.close()
        except Exception as exc:
            logger.error(f"Error extrayendo imágenes de {pdf_path.name}: {exc}")
        return extracted

    def extract_from_pptx(self, pptx_path: Path) -> list[tuple[int, Path]]:
        """Extrae imágenes de un PPTX con su índice de diapositiva."""
        extracted: list[tuple[int, Path]] = []
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            img = Image.open(io.BytesIO(shape.image.blob))
                            path = self.save_image(img)
                            if path:
                                extracted.append((slide_idx, path))
                        except Exception as exc:
                            logger.debug(f"Imagen en diapositiva {slide_idx}: {exc}")
        except Exception as exc:
            logger.error(f"Error extrayendo imágenes de {pptx_path.name}: {exc}")
        return extracted
