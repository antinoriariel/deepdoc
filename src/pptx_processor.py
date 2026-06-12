"""Procesamiento de presentaciones PPTX con extracción ordenada por posición visual."""
import io
from pathlib import Path
from typing import Optional

from PIL import Image
from loguru import logger

from .ocr_engine import OCREngine
from .image_extractor import ImageExtractor


class PPTXProcessor:
    def __init__(
        self,
        ocr_engine: OCREngine,
        image_extractor: Optional[ImageExtractor] = None,
        extract_images: bool = False,
        lang: Optional[list[str]] = None,
    ) -> None:
        self.ocr_engine = ocr_engine
        self.image_extractor = image_extractor
        self.extract_images = extract_images
        self.lang = lang or ["es", "en"]

    def process(self, pptx_path: Path) -> list[str]:
        """Procesa el PPTX y retorna bloques de texto ordenados."""
        blocks: list[str] = []
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            for slide_num, slide in enumerate(prs.slides):
                blocks.append(f"\n## Diapositiva {slide_num + 1}\n")
                blocks.extend(self._process_slide(slide, slide_num))
        except Exception as exc:
            logger.error(f"Error procesando {pptx_path.name}: {exc}")
        return blocks

    def get_slide_count(self, pptx_path: Path) -> int:
        """Retorna el número de diapositivas."""
        try:
            from pptx import Presentation

            return len(Presentation(str(pptx_path)).slides)
        except Exception:
            return 0

    # ------------------------------------------------------------------
    # Internos
    # ------------------------------------------------------------------

    def _process_slide(self, slide, slide_num: int) -> list[str]:
        """Extrae contenido de una diapositiva ordenado top-left → bottom-right."""
        blocks: list[str] = []
        try:
            sorted_shapes = sorted(
                slide.shapes,
                key=lambda s: (s.top if s.top is not None else 0, s.left if s.left is not None else 0),
            )
            for shape in sorted_shapes:
                blocks.extend(self._process_shape(shape, slide_num))

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    blocks.append(f"\n> **Notas:** {notes_text}\n")
        except Exception as exc:
            logger.warning(f"Error en diapositiva {slide_num + 1}: {exc}")
            blocks.append(f"\n> ⚠️ Diapositiva {slide_num + 1}: error al procesar.\n")
        return blocks

    def _process_shape(self, shape, slide_num: int) -> list[str]:
        """Despacha el procesamiento según el tipo de shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        blocks: list[str] = []
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                md = self._table_to_markdown(shape.table)
                if md:
                    blocks.append(md)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blocks.extend(self._process_image_shape(shape, slide_num))
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                blocks.extend(self._extract_text_frame(shape))
        except Exception as exc:
            logger.debug(f"Shape ignorado en diapositiva {slide_num + 1}: {exc}")
        return blocks

    def _extract_text_frame(self, shape) -> list[str]:
        """Extrae párrafos respetando niveles de lista y título."""
        blocks: list[str] = []
        is_title = any(kw in shape.name.lower() for kw in ("title", "título", "subtitle"))
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level
            if is_title and level == 0:
                blocks.append(f"### {text}")
            elif level > 0:
                indent = "  " * (level - 1)
                blocks.append(f"{indent}- {text}")
            else:
                blocks.append(text)
        return blocks

    def _table_to_markdown(self, table) -> str:
        """Convierte una tabla PPTX a formato Markdown."""
        rows = [
            [cell.text.strip().replace("\n", " ") for cell in row.cells]
            for row in table.rows
        ]
        if not rows:
            return ""
        col_count = len(rows[0])
        header = rows[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]
        for row in rows[1:]:
            padded = (row + [""] * col_count)[:col_count]
            lines.append("| " + " | ".join(padded) + " |")
        return "\n".join(lines)

    def _process_image_shape(self, shape, slide_num: int) -> list[str]:
        """Aplica OCR sobre una imagen de la diapositiva y la guarda si corresponde."""
        blocks: list[str] = []
        try:
            img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
            if self.extract_images and self.image_extractor:
                path = self.image_extractor.save_image(img)
                if path:
                    blocks.append(f"\n![Figura](images/{path.name})\n")
            ocr_text = self.ocr_engine.extract_text(img, self.lang)
            if ocr_text.strip():
                blocks.append(ocr_text)
        except Exception as exc:
            logger.debug(f"Imagen ignorada en diapositiva {slide_num + 1}: {exc}")
        return blocks
